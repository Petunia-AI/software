"""
Presentations API — Generate AI presentations from meeting transcripts using Claude.

Endpoints:
  GET    /presentations               — List presentations (most recent first)
  POST   /presentations               — Create + generate presentation from transcript
  GET    /presentations/{id}          — Get full presentation detail (includes HTML)
  DELETE /presentations/{id}          — Delete presentation
  POST   /presentations/{id}/regenerate — Regenerate with (optionally different) style
  GET    /presentations/{id}/download.pptx — Download as editable PowerPoint
"""
from __future__ import annotations

import re
from datetime import datetime, timezone
from typing import Optional

import structlog
from anthropic import AsyncAnthropic
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.api.auth import get_current_user
from app.config import settings
from app.database import get_db
from app.models.business import Business
from app.models.presentation import Presentation
from app.models.user import User

logger = structlog.get_logger()
router = APIRouter(prefix="/presentations", tags=["presentations"])

# ── Style configurations ───────────────────────────────────────────────────────

STYLE_CONFIGS: dict[str, dict] = {
    "profesional": {
        "primary": "#7C3AED",
        "secondary": "#5B21B6",
        "accent": "#C4B5FD",
        "bg": "#FFFFFF",
        "desc": "Elegante con gradientes violeta y púrpura. Tipografía moderna y limpia.",
    },
    "creativo": {
        "primary": "#EC4899",
        "secondary": "#8B5CF6",
        "accent": "#F59E0B",
        "bg": "#FAFAFA",
        "desc": "Vibrante y dinámico. Colores llamativos, formas modernas.",
    },
    "minimalista": {
        "primary": "#111827",
        "secondary": "#374151",
        "accent": "#6B7280",
        "bg": "#FFFFFF",
        "desc": "Limpio y sofisticado. Mucho espacio en blanco, tipografía grande.",
    },
    "corporativo": {
        "primary": "#1E3A5F",
        "secondary": "#2563EB",
        "accent": "#93C5FD",
        "bg": "#F8FAFC",
        "desc": "Formal y estructurado. Azul marino, ideal para presentaciones ejecutivas.",
    },
}

# ── Helpers ────────────────────────────────────────────────────────────────────

async def _get_business(db: AsyncSession, user: User) -> Business:
    r = await db.execute(select(Business).where(Business.id == user.business_id))
    biz = r.scalar_one_or_none()
    if not biz:
        raise HTTPException(status_code=404, detail="Negocio no encontrado")
    return biz


async def _get_presentation(db: AsyncSession, pres_id: str, business_id: str) -> Presentation:
    r = await db.execute(
        select(Presentation).where(
            Presentation.id == pres_id,
            Presentation.business_id == business_id,
        )
    )
    p = r.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Presentación no encontrada")
    return p


async def _generate_html(
    title: str,
    transcript: str,
    style: str,
    business_name: str,
) -> tuple[str, int]:
    """Call Claude to generate a beautiful HTML presentation. Returns (html, slide_count)."""

    cfg = STYLE_CONFIGS.get(style, STYLE_CONFIGS["profesional"])
    today = datetime.now(timezone.utc).strftime("%d de %B de %Y")

    if len(transcript) > 14000:
        transcript = transcript[:14000] + "\n\n[... transcripción truncada ...]"

    prompt = f"""Eres un diseñador experto en presentaciones ejecutivas de nivel Fortune 500. Genera una presentación HTML espectacular, moderna y lista para usar.

DATOS:
- Título: {title}
- Empresa: {business_name}
- Fecha: {today}
- Estilo: {style} — {cfg["desc"]}
- Color primario: {cfg["primary"]} | Secundario: {cfg["secondary"]} | Acento: {cfg["accent"]}

TRANSCRIPCIÓN:
{transcript}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
GENERA UN DOCUMENTO HTML COMPLETO, auto-contenido y espectacular:

El HTML debe:
1. Comenzar con <!DOCTYPE html> y terminar con </html>
2. Incluir Google Fonts: Inter (para cuerpo) + optionally Playfair Display (para títulos)
3. Tener CSS completo en <style> dentro del <head>
4. Tener las diapositivas como divs con class="slide"
5. NO usar JavaScript
6. NO usar imágenes externas (solo CSS shapes)

CSS OBLIGATORIO (cópialo exactamente, no lo omitas):
<style>
*, *::before, *::after {{ margin: 0; padding: 0; box-sizing: border-box; }}
html, body {{ background: #E8E5F0; font-family: 'Inter', 'Segoe UI', sans-serif; }}
.slide {{
  width: 1280px; height: 720px;
  margin: 32px auto; overflow: hidden;
  position: relative; border-radius: 20px;
  box-shadow: 0 20px 60px rgba(0,0,0,0.18), 0 4px 16px rgba(0,0,0,0.08);
  display: flex; flex-direction: column;
  background: #fff;
}}
@media print {{
  html, body {{ background: white; }}
  .slide {{
    width: 297mm; height: 167mm;
    margin: 0; border-radius: 0;
    box-shadow: none;
    page-break-after: always; page-break-inside: avoid;
  }}
}}
</style>

DISEÑO POR TIPO DE SLIDE:

SLIDE 1 — PORTADA (espectacular):
```html
<div class="slide" style="background: linear-gradient(135deg, {cfg["primary"]} 0%, {cfg["secondary"]} 60%, #1a0a3e 100%);">
  <!-- Círculos decorativos -->
  <div style="position:absolute; width:500px; height:500px; border-radius:50%; background: rgba(255,255,255,0.05); top:-180px; right:-120px;"></div>
  <div style="position:absolute; width:300px; height:300px; border-radius:50%; background: rgba(255,255,255,0.04); bottom:-100px; left:-80px;"></div>
  <div style="position:absolute; width:180px; height:180px; border-radius:50%; border: 2px solid rgba(255,255,255,0.15); top:40px; right:80px;"></div>
  <!-- Barra lateral de acento -->
  <div style="position:absolute; left:0; top:0; width:6px; height:100%; background:{cfg["accent"]};"></div>
  <!-- Contenido centrado -->
  <div style="flex:1; display:flex; flex-direction:column; justify-content:center; padding: 60px 80px;">
    <div style="font-size:13px; font-weight:600; letter-spacing:3px; color:{cfg["accent"]}; text-transform:uppercase; margin-bottom:24px;">{business_name}</div>
    <h1 style="font-size:58px; font-weight:800; color:#fff; line-height:1.1; max-width:800px; margin-bottom:28px;">{title}</h1>
    <div style="width:80px; height:4px; background:{cfg["accent"]}; border-radius:2px; margin-bottom:28px;"></div>
    <div style="font-size:20px; color:rgba(255,255,255,0.65);">{today}</div>
  </div>
  <!-- Footer -->
  <div style="background:rgba(0,0,0,0.25); padding:16px 80px; font-size:12px; color:rgba(255,255,255,0.4); display:flex; justify-content:space-between;">
    <span>Confidencial</span><span>{business_name} © 2026</span>
  </div>
</div>
```

SLIDES DE CONTENIDO (usa este patrón):
```html
<div class="slide" style="background:{cfg["bg"]};">
  <!-- Header con gradiente -->
  <div style="background: linear-gradient(90deg, {cfg["primary"]}, {cfg["secondary"]}); padding: 0 60px; height:110px; display:flex; align-items:center; gap:20px; flex-shrink:0;">
    <div style="width:4px; height:50px; background:{cfg["accent"]}; border-radius:2px; flex-shrink:0;"></div>
    <h2 style="font-size:28px; font-weight:700; color:#fff; letter-spacing:-0.3px;">TÍTULO DEL SLIDE</h2>
    <div style="margin-left:auto; font-size:12px; color:rgba(255,255,255,0.4); font-weight:500;">N / TOTAL</div>
  </div>
  <!-- Línea decorativa -->
  <div style="height:4px; background: linear-gradient(90deg, {cfg["accent"]}, transparent);"></div>
  <!-- Contenido con bullets -->
  <div style="flex:1; padding: 40px 60px; display:flex; flex-direction:column; gap:18px;">
    <div style="display:flex; align-items:flex-start; gap:16px;">
      <div style="width:8px; height:8px; border-radius:50%; background:{cfg["primary"]}; margin-top:8px; flex-shrink:0;"></div>
      <p style="font-size:20px; color:#1a1a2e; line-height:1.5; font-weight:400;">Punto clave aquí</p>
    </div>
  </div>
  <!-- Footer -->
  <div style="border-top: 1px solid #f0eef8; padding: 12px 60px; display:flex; justify-content:space-between; align-items:center; flex-shrink:0;">
    <span style="font-size:11px; color:#999; font-weight:500;">{business_name}</span>
    <div style="width:40px; height:3px; background:{cfg["accent"]}; border-radius:2px;"></div>
    <span style="font-size:11px; color:#999;">{today}</span>
  </div>
</div>
```

SLIDE DE MÉTRICAS (si hay datos numéricos):
- Tarjetas en grid de 3-4 con número grande (60px, primary), etiqueta descriptiva (16px, gray)
- Cada tarjeta: fondo suave (#F5F3FF o similar), borde top de 4px con primary color
- Los números deben venir de la transcripción real

SLIDE DE CIERRE — PRÓXIMOS PASOS:
- Fondo oscuro: `background: linear-gradient(135deg, #0f0f1a, {cfg["secondary"]})`
- Lista con ✓ en color acento
- Texto blanco

ESTRUCTURA EXACTA DE SLIDES (8-12 según contenido):
1. Portada espectacular
2. Agenda / Índice de temas
3-9. Contenido extraído de la transcripción (tema importante por slide)
10. Métricas y datos (solo si hay números en la transcripción)
11. Compromisos y acuerdos
12. Próximos pasos y cierre

REGLAS DE ORO:
- Máximo 5-6 puntos por slide, NUNCA párrafos largos
- TODO EN ESPAÑOL
- Extrae información real de la transcripción (no inventes)
- El número de slide va en el header (ej: "3 / 10")
- Bullets con iconos ▸ o • o números según el tipo de slide
- Datos numéricos destacados visualmente

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Responde ÚNICAMENTE con el HTML completo. Sin explicaciones. Sin markdown. Sin ``` code blocks.
Empieza con <!DOCTYPE html> y termina con </html>."""

    client = AsyncAnthropic(api_key=settings.anthropic_api_key)
    message = await client.messages.create(
        model=settings.claude_model,
        max_tokens=8000,
        messages=[{"role": "user", "content": prompt}],
    )

    html = message.content[0].text.strip()

    # Strip markdown code fences if Claude added them
    if html.startswith("```"):
        html = re.sub(r"^```[a-z]*\n?", "", html)
        html = re.sub(r"\n?```$", "", html)
        html = html.strip()

    slide_count = len(re.findall(r'class=["\'][^"\']*slide[^"\']*["\']', html))

    return html, max(slide_count, 1)


async def _generate_pptx(
    title: str,
    transcript: str,
    style: str,
    business_name: str,
) -> bytes:
    """Ask Claude for structured slide content, then build a beautiful PPTX."""
    import json as _json
    import io
    from pptx import Presentation as PPTXPres
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    cfg = STYLE_CONFIGS.get(style, STYLE_CONFIGS["profesional"])
    today = datetime.now(timezone.utc).strftime("%d de %B de %Y")

    if len(transcript) > 12000:
        transcript = transcript[:12000] + "\n\n[... transcripción truncada ...]"

    client = AsyncAnthropic(api_key=settings.anthropic_api_key)
    prompt = f"""Analiza la siguiente transcripción y genera el contenido estructurado para una presentación ejecutiva.

DATOS:
- Título: {title}
- Empresa: {business_name}
- Fecha: {today}

TRANSCRIPCIÓN:
{transcript}

Responde ÚNICAMENTE con JSON válido (sin markdown, sin triple backticks), exactamente así:
{{
  "slides": [
    {{"type": "cover", "title": "{title}", "subtitle": "{business_name} · {today}"}},
    {{"type": "content", "title": "Agenda", "bullets": ["Tema tratado 1", "Tema tratado 2", "Tema tratado 3"]}},
    {{"type": "content", "title": "Título del tema principal", "bullets": ["Punto concreto 1", "Dato o acuerdo relevante", "Insight clave de la reunión"]}},
    {{"type": "metrics", "title": "Cifras y Resultados", "metrics": [{{"value": "23%", "label": "Crecimiento Q1"}}, {{"value": "$500K", "label": "Meta Revenue"}}]}},
    {{"type": "closing", "title": "Próximos Pasos", "bullets": ["Acción 1 con responsable o fecha si la hay", "Acción 2", "Acción 3"]}}
  ]
}}

REGLAS:
- 8-12 slides total
- Solo 1 slide cover y 1 closing
- Incluir metrics solo si hay datos numéricos reales en la transcripción
- Máximo 5 bullets por slide de contenido
- Todo en español, conciso (máx 90 chars por bullet)
- Extraer información real de la transcripción"""

    message = await client.messages.create(
        model=settings.claude_model,
        max_tokens=3000,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = message.content[0].text.strip()
    if "```" in raw:
        raw = re.sub(r"^```[a-z]*\n?", "", raw)
        raw = re.sub(r"\n?```$", "", raw)
        raw = raw.strip()

    slides_data = _json.loads(raw).get("slides", [])

    # ── Build the PPTX ────────────────────────────────────────────────────────
    prs = PPTXPres()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def hex_rgb(h: str) -> RGBColor:
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    PRIMARY   = hex_rgb(cfg["primary"])
    SECONDARY = hex_rgb(cfg["secondary"])
    ACCENT    = hex_rgb(cfg["accent"])
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    DARK      = RGBColor(0x0F, 0x0F, 0x1A)
    LIGHT_BG  = RGBColor(0xF5, 0xF3, 0xFF)
    GRAY      = RGBColor(0x64, 0x74, 0x8B)

    W = prs.slide_width
    H = prs.slide_height

    def add_rect(slide, x, y, w, h, fill, line_rgb=None):
        s = slide.shapes.add_shape(1, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if line_rgb:
            s.line.color.rgb = line_rgb
            s.line.width = Pt(1.5)
        else:
            s.line.fill.background()
        return s

    def add_oval(slide, x, y, w, h, fill):
        s = slide.shapes.add_shape(9, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.fill.background()
        return s

    def set_text(tf, text, size, bold=False, color=None, align=None, space_before=0):
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        if align:
            p.alignment = align
        if space_before:
            p.space_before = Pt(space_before)
        tf.word_wrap = True

    slide_total = len(slides_data)

    for idx, sd in enumerate(slides_data):
        stype = sd.get("type", "content")
        slide = prs.slides.add_slide(blank)

        if stype == "cover":
            # Background
            add_rect(slide, 0, 0, W, H, PRIMARY)
            # Decorative circles
            add_oval(slide, Inches(10.8), Inches(-1.5), Inches(5), Inches(5), SECONDARY)
            add_oval(slide, Inches(-1.2), Inches(5.5), Inches(3.5), Inches(3.5), SECONDARY)
            add_oval(slide, Inches(9.5), Inches(5), Inches(2), Inches(2), ACCENT)
            # Left accent bar
            add_rect(slide, 0, 0, Inches(0.09), H, ACCENT)
            # Bottom bar
            add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55), SECONDARY)

            # Business name tag
            tag = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(10), Inches(0.5))
            set_text(tag.text_frame, business_name.upper(), 11, bold=True, color=ACCENT)

            # Main title
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(10.5), Inches(2.8))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]
            p.text = sd.get("title", title)
            p.font.size = Pt(50)
            p.font.bold = True
            p.font.color.rgb = WHITE

            # Accent line
            add_rect(slide, Inches(1.0), Inches(5.2), Inches(1.2), Inches(0.07), ACCENT)

            # Subtitle / date
            sub = slide.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(10), Inches(0.6))
            set_text(sub.text_frame, sd.get("subtitle", f"{business_name} · {today}"), 18, color=RGBColor(0xCC, 0xBB, 0xFF))

            # Footer text
            ft = slide.shapes.add_textbox(Inches(0.5), H - Inches(0.45), Inches(6), Inches(0.35))
            set_text(ft.text_frame, "Confidencial · Uso interno", 10, color=RGBColor(0x99, 0x88, 0xCC))

        elif stype in ("content",):
            # Header bar
            add_rect(slide, 0, 0, W, Inches(1.2), PRIMARY)
            add_rect(slide, 0, Inches(1.2), W, Inches(0.07), ACCENT)
            # Accent left square in header
            add_rect(slide, Inches(0.5), Inches(0.25), Inches(0.06), Inches(0.7), ACCENT)

            # Title
            th = slide.shapes.add_textbox(Inches(0.75), Inches(0.18), Inches(11.2), Inches(0.85))
            set_text(th.text_frame, sd.get("title", ""), 27, bold=True, color=WHITE)

            # Slide number
            sn = slide.shapes.add_textbox(Inches(12.0), Inches(0.38), Inches(1.0), Inches(0.45))
            set_text(sn.text_frame, f"{idx+1} / {slide_total}", 11, color=RGBColor(0xCC, 0xBB, 0xFF), align=PP_ALIGN.RIGHT)

            # Bullets
            bullets = sd.get("bullets", [])
            if bullets:
                cb = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.5), Inches(5.5))
                tf = cb.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(bullets[:6]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    r1 = p.add_run()
                    r1.text = "▸  "
                    r1.font.color.rgb = PRIMARY
                    r1.font.size = Pt(16)
                    r1.font.bold = True
                    r2 = p.add_run()
                    r2.text = bullet
                    r2.font.size = Pt(20)
                    r2.font.color.rgb = DARK
                    p.space_before = Pt(10)
                    p.space_after = Pt(4)

            # Bottom footer line
            add_rect(slide, 0, H - Inches(0.22), W, Inches(0.22), PRIMARY)
            ft = slide.shapes.add_textbox(Inches(0.5), H - Inches(0.2), Inches(6), Inches(0.18))
            set_text(ft.text_frame, business_name, 9, color=WHITE)

        elif stype == "metrics":
            # Header
            add_rect(slide, 0, 0, W, Inches(1.2), PRIMARY)
            add_rect(slide, 0, Inches(1.2), W, Inches(0.07), ACCENT)
            add_rect(slide, Inches(0.5), Inches(0.25), Inches(0.06), Inches(0.7), ACCENT)

            th = slide.shapes.add_textbox(Inches(0.75), Inches(0.18), Inches(11.2), Inches(0.85))
            set_text(th.text_frame, sd.get("title", "Cifras Clave"), 27, bold=True, color=WHITE)

            sn = slide.shapes.add_textbox(Inches(12.0), Inches(0.38), Inches(1.0), Inches(0.45))
            set_text(sn.text_frame, f"{idx+1} / {slide_total}", 11, color=RGBColor(0xCC, 0xBB, 0xFF), align=PP_ALIGN.RIGHT)

            metrics = sd.get("metrics", [])[:4]
            n = max(len(metrics), 1)
            gap = Inches(0.35)
            card_w = (W - gap * (n + 1)) / n
            card_h = Inches(3.9)
            card_y = Inches(1.65)

            for i, m in enumerate(metrics):
                cx = gap + i * (card_w + gap)
                # Card shadow effect (darker rect behind)
                shadow = slide.shapes.add_shape(1, cx + Inches(0.05), card_y + Inches(0.05), card_w, card_h)
                shadow.fill.solid()
                shadow.fill.fore_color.rgb = RGBColor(0xD0, 0xC8, 0xF0)
                shadow.line.fill.background()

                card = add_rect(slide, cx, card_y, card_w, card_h, LIGHT_BG, line_rgb=ACCENT)
                # Top accent strip
                add_rect(slide, cx, card_y, card_w, Inches(0.22), PRIMARY)

                val_box = slide.shapes.add_textbox(cx + Inches(0.1), card_y + Inches(0.4), card_w - Inches(0.2), Inches(1.9))
                val_box.text_frame.word_wrap = True
                pv = val_box.text_frame.paragraphs[0]
                pv.text = m.get("value", "—")
                pv.font.size = Pt(48)
                pv.font.bold = True
                pv.font.color.rgb = PRIMARY
                pv.alignment = PP_ALIGN.CENTER

                lbl_box = slide.shapes.add_textbox(cx + Inches(0.1), card_y + Inches(2.6), card_w - Inches(0.2), Inches(1.1))
                lbl_box.text_frame.word_wrap = True
                pl = lbl_box.text_frame.paragraphs[0]
                pl.text = m.get("label", "")
                pl.font.size = Pt(16)
                pl.font.color.rgb = GRAY
                pl.alignment = PP_ALIGN.CENTER

            add_rect(slide, 0, H - Inches(0.22), W, Inches(0.22), PRIMARY)

        elif stype == "closing":
            # Dark background
            add_rect(slide, 0, 0, W, H, DARK)
            # Decorative circles
            add_oval(slide, Inches(10.5), Inches(3.8), Inches(4.5), Inches(4.5), PRIMARY)
            add_oval(slide, Inches(-1.5), Inches(-0.8), Inches(4), Inches(4), SECONDARY)
            # Top accent bar
            add_rect(slide, 0, 0, W, Inches(0.5), PRIMARY)
            # Left bar
            add_rect(slide, 0, 0, Inches(0.09), H, ACCENT)

            th = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(1.1))
            set_text(th.text_frame, sd.get("title", "Próximos Pasos"), 36, bold=True, color=WHITE)

            # Divider
            add_rect(slide, Inches(0.8), Inches(1.85), Inches(5), Inches(0.06), ACCENT)

            bullets = sd.get("bullets", [])
            if bullets:
                cb = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.5))
                tf = cb.text_frame
                tf.word_wrap = True
                for i, b in enumerate(bullets[:5]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    r1 = p.add_run()
                    r1.text = "✓  "
                    r1.font.color.rgb = ACCENT
                    r1.font.size = Pt(18)
                    r1.font.bold = True
                    r2 = p.add_run()
                    r2.text = b
                    r2.font.size = Pt(20)
                    r2.font.color.rgb = RGBColor(0xE2, 0xD9, 0xFF)
                    p.space_before = Pt(14)

            # "Gracias" at bottom
            thanks = slide.shapes.add_textbox(0, H - Inches(1.1), W, Inches(0.7))
            set_text(thanks.text_frame, "¿Preguntas?", 26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Schemas ────────────────────────────────────────────────────────────────────

class CreatePresentationRequest(BaseModel):
    title: str
    transcript_text: str
    style: str = "profesional"


class RegenerateRequest(BaseModel):
    style: Optional[str] = None


class PresentationListItem(BaseModel):
    id: str
    title: str
    style: str
    slide_count: int
    created_at: datetime
    model_config = {"from_attributes": True}


class PresentationOut(BaseModel):
    id: str
    business_id: str
    title: str
    transcript_text: Optional[str]
    style: str
    presentation_html: Optional[str]
    slide_count: int
    created_at: datetime
    updated_at: datetime
    model_config = {"from_attributes": True}


# ── Endpoints ──────────────────────────────────────────────────────────────────

@router.get("", response_model=list[PresentationListItem])
async def list_presentations(
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)
    r = await db.execute(
        select(Presentation)
        .where(Presentation.business_id == biz.id)
        .order_by(Presentation.created_at.desc())
    )
    return r.scalars().all()


@router.post("", response_model=PresentationOut)
async def create_presentation(
    payload: CreatePresentationRequest,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)

    if not payload.transcript_text or len(payload.transcript_text.strip()) < 50:
        raise HTTPException(status_code=400, detail="La transcripción debe tener al menos 50 caracteres.")
    if not settings.anthropic_api_key:
        raise HTTPException(status_code=503, detail="ANTHROPIC_API_KEY no configurada.")

    style = payload.style if payload.style in STYLE_CONFIGS else "profesional"
    logger.info("generating_presentation", title=payload.title, style=style, business=biz.id)

    html, slide_count = await _generate_html(
        title=payload.title,
        transcript=payload.transcript_text,
        style=style,
        business_name=biz.name or "Empresa",
    )

    pres = Presentation(
        business_id=biz.id,
        title=payload.title,
        transcript_text=payload.transcript_text,
        style=style,
        presentation_html=html,
        slide_count=slide_count,
    )
    db.add(pres)
    await db.commit()
    await db.refresh(pres)

    logger.info("presentation_created", id=pres.id, slides=slide_count)
    return pres


@router.get("/{presentation_id}", response_model=PresentationOut)
async def get_presentation(
    presentation_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)
    return await _get_presentation(db, presentation_id, biz.id)


@router.delete("/{presentation_id}", status_code=204)
async def delete_presentation(
    presentation_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)
    pres = await _get_presentation(db, presentation_id, biz.id)
    await db.delete(pres)
    await db.commit()


@router.post("/{presentation_id}/regenerate", response_model=PresentationOut)
async def regenerate_presentation(
    presentation_id: str,
    payload: RegenerateRequest,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)
    pres = await _get_presentation(db, presentation_id, biz.id)

    if not pres.transcript_text:
        raise HTTPException(status_code=400, detail="No hay transcripción guardada para regenerar.")
    if not settings.anthropic_api_key:
        raise HTTPException(status_code=503, detail="ANTHROPIC_API_KEY no configurada.")

    style = payload.style if payload.style in STYLE_CONFIGS else pres.style

    html, slide_count = await _generate_html(
        title=pres.title,
        transcript=pres.transcript_text,
        style=style,
        business_name=biz.name or "Empresa",
    )

    pres.presentation_html = html
    pres.slide_count = slide_count
    pres.style = style
    await db.commit()
    await db.refresh(pres)

    logger.info("presentation_regenerated", id=pres.id, slides=slide_count, style=style)
    return pres


@router.get("/{presentation_id}/download.pptx")
async def download_pptx(
    presentation_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """Generate and download the presentation as an editable PowerPoint (.pptx) file."""
    biz = await _get_business(db, current_user)
    pres = await _get_presentation(db, presentation_id, biz.id)

    if not pres.transcript_text:
        raise HTTPException(status_code=400, detail="Sin transcripción para generar el PowerPoint.")
    if not settings.anthropic_api_key:
        raise HTTPException(status_code=503, detail="ANTHROPIC_API_KEY no configurada.")

    try:
        pptx_bytes = await _generate_pptx(
            title=pres.title,
            transcript=pres.transcript_text,
            style=pres.style,
            business_name=biz.name or "Empresa",
        )
    except Exception as e:
        logger.error("pptx_generation_error", id=presentation_id, error=str(e))
        raise HTTPException(status_code=500, detail=f"Error al generar el PowerPoint: {e}")

    safe_name = re.sub(r"[^\w\s-]", "", pres.title).strip().replace(" ", "_")[:60]
    filename = f"{safe_name}_presentacion.pptx"

    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )

import re
from datetime import datetime, timezone
from typing import Optional

import structlog
from anthropic import AsyncAnthropic
from fastapi import APIRouter, Depends, HTTPException
from pydantic import BaseModel
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.api.auth import get_current_user
from app.config import settings
from app.database import get_db
from app.models.business import Business
from app.models.presentation import Presentation
from app.models.user import User

logger = structlog.get_logger()
router = APIRouter(prefix="/presentations", tags=["presentations"])

# ── Style configurations ───────────────────────────────────────────────────────

STYLE_CONFIGS: dict[str, dict] = {
    "profesional": {
        "primary": "#7C3AED",
        "secondary": "#5B21B6",
        "accent": "#C4B5FD",
        "bg": "#FFFFFF",
        "desc": "Elegante con gradientes violeta y púrpura. Tipografía moderna y limpia.",
    },
    "creativo": {
        "primary": "#EC4899",
        "secondary": "#8B5CF6",
        "accent": "#F59E0B",
        "bg": "#FAFAFA",
        "desc": "Vibrante y dinámico. Colores llamativos, formas modernas.",
    },
    "minimalista": {
        "primary": "#111827",
        "secondary": "#374151",
        "accent": "#6B7280",
        "bg": "#FFFFFF",
        "desc": "Limpio y sofisticado. Mucho espacio en blanco, tipografía grande.",
    },
    "corporativo": {
        "primary": "#1E3A5F",
        "secondary": "#2563EB",
        "accent": "#93C5FD",
        "bg": "#F8FAFC",
        "desc": "Formal y estructurado. Azul marino, ideal para presentaciones ejecutivas.",
    },
}

# ── Helpers ────────────────────────────────────────────────────────────────────

async def _get_business(db: AsyncSession, user: User) -> Business:
    r = await db.execute(select(Business).where(Business.id == user.business_id))
    biz = r.scalar_one_or_none()
    if not biz:
        raise HTTPException(status_code=404, detail="Negocio no encontrado")
    return biz


async def _get_presentation(db: AsyncSession, pres_id: str, business_id: str) -> Presentation:
    r = await db.execute(
        select(Presentation).where(
            Presentation.id == pres_id,
            Presentation.business_id == business_id,
        )
    )
    p = r.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Presentación no encontrada")
    return p


async def _generate_html(
    title: str,
    transcript: str,
    style: str,
    business_name: str,
) -> tuple[str, int]:
    """Call Claude to generate a beautiful HTML presentation. Returns (html, slide_count)."""

    cfg = STYLE_CONFIGS.get(style, STYLE_CONFIGS["profesional"])
    today = datetime.now(timezone.utc).strftime("%d de %B de %Y")

    # Truncate very long transcripts to avoid token limits
    if len(transcript) > 14000:
        transcript = transcript[:14000] + "\n\n[... transcripción truncada ...]"

    prompt = f"""Eres un diseñador experto en presentaciones ejecutivas. Genera una presentación HTML espectacular y lista para usar.

DATOS DE LA PRESENTACIÓN:
- Título: {title}
- Empresa: {business_name}
- Fecha: {today}
- Estilo: {style} — {cfg["desc"]}
- Color primario: {cfg["primary"]} | Secundario: {cfg["secondary"]} | Acento: {cfg["accent"]}

TRANSCRIPCIÓN DE LA REUNIÓN:
{transcript}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
GENERA UN DOCUMENTO HTML COMPLETO con estas especificaciones exactas:

## ESTRUCTURA HTML:
```
<!DOCTYPE html>
<html lang="es">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>[título]</title>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap" rel="stylesheet">
  <style>
    /* ... todo el CSS aquí ... */
  </style>
</head>
<body>
  <div class="slide"> ... </div>
  <div class="slide"> ... </div>
  ...
</body>
</html>
```

## CSS OBLIGATORIO:
```css
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
html, body {{ background: #F1F5F9; font-family: 'Inter', sans-serif; }}
.slide {{
  width: 1280px;
  height: 720px;
  margin: 24px auto;
  overflow: hidden;
  position: relative;
  border-radius: 16px;
  box-shadow: 0 8px 40px rgba(0,0,0,0.12);
  display: flex;
  flex-direction: column;
}}
@media print {{
  body {{ background: white; }}
  .slide {{
    width: 100%; height: 100vh;
    margin: 0; border-radius: 0;
    box-shadow: none; page-break-after: always;
  }}
}}
```

## DISEÑO POR TIPO DE SLIDE:

### SLIDE 1 — PORTADA:
- Fondo: gradiente espectacular con {cfg["primary"]}
- Título: 64px, blanco, centrado, font-weight 800
- Subtítulo: nombre empresa + fecha en color semi-transparente
- Elemento decorativo: círculos o formas geométricas semitransparentes en el fondo

### SLIDES DE CONTENIDO (estándar):
- Header: banda de 110px con fondo {cfg["primary"]}, título del slide en blanco 28px font-weight 700
- Body: fondo {cfg["bg"]}, padding 48px 60px
- Bullets: puntos con color {cfg["primary"]}, texto 18px, line-height 2
- Máximo 5-6 puntos por slide, nunca párrafos largos

### SLIDE DE MÉTRICAS (si hay datos numéricos en la transcripción):
- Grid de 3-4 tarjetas con: número grande (72px bold, color {cfg["primary"]}), etiqueta descriptiva
- Fondo de cada tarjeta: gradiente suave

### SLIDE DE COMPROMISOS/ACUERDOS:
- Lista con iconos ✓ o ▸ en verde
- Fondo oscuro tipo pizarrón o gradiente oscuro con texto claro

### SLIDE FINAL — PRÓXIMOS PASOS Y CIERRE:
- Fondo oscuro con gradiente
- Lista de acciones con fechas si las hay
- "Gracias" o "¿Preguntas?" al final

## ESTRUCTURA DE SLIDES (8-12 slides según el contenido):
1. Portada con título y subtítulo
2. Agenda / Temas tratados hoy
3-8. Contenido principal extraído de la transcripción (un tema importante por slide)
9. Compromisos y acuerdos clave
10. Próximos pasos
11. (Opcional) Datos/métricas si hay números importantes
12. Cierre / Gracias

## REGLAS DE ORO:
- Texto MÍNIMO por slide — solo puntos clave, nunca párrafos
- Extrae y destaca DATOS NUMÉRICOS (fechas, porcentajes, montos, plazos)
- TODO EN ESPAÑOL siempre, sin importar el idioma de la transcripción
- Sin JavaScript — HTML y CSS puros únicamente
- El HTML debe ser completamente auto-contenido
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Responde ÚNICAMENTE con el HTML completo. Sin explicaciones. Sin markdown. Sin ``` code blocks.
Comienza directamente con <!DOCTYPE html> y termina con </html>."""

    client = AsyncAnthropic(api_key=settings.anthropic_api_key)
    message = await client.messages.create(
        model="claude-opus-4-5",
        max_tokens=8000,
        messages=[{"role": "user", "content": prompt}],
    )

    html = message.content[0].text.strip()

    # Strip markdown code fences if Claude added them
    if html.startswith("```"):
        html = re.sub(r"^```[a-z]*\n?", "", html)
        html = re.sub(r"\n?```$", "", html)
        html = html.strip()

    # Count slides
    slide_count = len(re.findall(r'class=["\']slide["\']', html))

    return html, slide_count


# ── Schemas ────────────────────────────────────────────────────────────────────

class CreatePresentationRequest(BaseModel):
    title: str
    transcript_text: str
    style: str = "profesional"


class RegenerateRequest(BaseModel):
    style: Optional[str] = None


class PresentationListItem(BaseModel):
    id: str
    title: str
    style: str
    slide_count: int
    created_at: datetime
    model_config = {"from_attributes": True}


class PresentationOut(BaseModel):
    id: str
    business_id: str
    title: str
    transcript_text: Optional[str]
    style: str
    presentation_html: Optional[str]
    slide_count: int
    created_at: datetime
    updated_at: datetime
    model_config = {"from_attributes": True}


# ── Endpoints ──────────────────────────────────────────────────────────────────

@router.get("", response_model=list[PresentationListItem])
async def list_presentations(
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)
    r = await db.execute(
        select(Presentation)
        .where(Presentation.business_id == biz.id)
        .order_by(Presentation.created_at.desc())
    )
    return r.scalars().all()


@router.post("", response_model=PresentationOut)
async def create_presentation(
    payload: CreatePresentationRequest,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)

    if not payload.transcript_text or len(payload.transcript_text.strip()) < 50:
        raise HTTPException(
            status_code=400,
            detail="La transcripción debe tener al menos 50 caracteres.",
        )
    if not settings.anthropic_api_key:
        raise HTTPException(status_code=503, detail="ANTHROPIC_API_KEY no configurada.")

    style = payload.style if payload.style in STYLE_CONFIGS else "profesional"

    logger.info("generating_presentation", title=payload.title, style=style, business=biz.id)

    html, slide_count = await _generate_html(
        title=payload.title,
        transcript=payload.transcript_text,
        style=style,
        business_name=biz.name or "Empresa",
    )

    pres = Presentation(
        business_id=biz.id,
        title=payload.title,
        transcript_text=payload.transcript_text,
        style=style,
        presentation_html=html,
        slide_count=slide_count,
    )
    db.add(pres)
    await db.commit()
    await db.refresh(pres)

    logger.info("presentation_created", id=pres.id, slides=slide_count)
    return pres


@router.get("/{presentation_id}", response_model=PresentationOut)
async def get_presentation(
    presentation_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)
    return await _get_presentation(db, presentation_id, biz.id)


@router.delete("/{presentation_id}", status_code=204)
async def delete_presentation(
    presentation_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)
    pres = await _get_presentation(db, presentation_id, biz.id)
    await db.delete(pres)
    await db.commit()


@router.post("/{presentation_id}/regenerate", response_model=PresentationOut)
async def regenerate_presentation(
    presentation_id: str,
    payload: RegenerateRequest,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    biz = await _get_business(db, current_user)
    pres = await _get_presentation(db, presentation_id, biz.id)

    if not pres.transcript_text:
        raise HTTPException(
            status_code=400, detail="No hay transcripción guardada para regenerar."
        )

    style = payload.style if payload.style in STYLE_CONFIGS else pres.style

    html, slide_count = await _generate_html(
        title=pres.title,
        transcript=pres.transcript_text,
        style=style,
        business_name=biz.name or "Empresa",
    )

    pres.presentation_html = html
    pres.slide_count = slide_count
    pres.style = style
    await db.commit()
    await db.refresh(pres)

    logger.info("presentation_regenerated", id=pres.id, slides=slide_count, style=style)
    return pres
